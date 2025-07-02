import streamlit as st
import pandas as pd
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

st.set_page_config(page_title="Générateur de slide", layout="centered")
st.title("Générateur de slide PowerPoint")

uploaded = st.file_uploader("Déposez votre fichier Excel (.xlsx)", type="xlsx")
if not uploaded:
    st.info("En attente de votre fichier…")
else:
    # lecture & transformation
    df = pd.read_excel(uploaded)
    df['Mois'] = pd.to_datetime(df['Mois'])
    df = df.sort_values('Mois')
    df['Durée'] = (
        df['Temps effectif total']
          .str.split(':')
          .apply(lambda x: int(x[0]) + int(x[1]) / 60)
    )
    pivot = (
        df
        .pivot_table(
            index='Mois',
            columns='E-formation',
            values='Durée',
            aggfunc='sum'
        )
        .reindex(sorted(df['Mois'].unique()), fill_value=0)
        .cumsum()
        .reset_index()
    )
    pivot['Mois_str'] = pivot['Mois'].dt.strftime('%b %Y')

    # création du graphique
    fig = go.Figure()
    fig.add_trace(go.Scatter(
        x=pivot['Mois_str'],
        y=pivot.get('Bibliothèque Numérique ENI', []),
        mode='lines+markers',
        name='Bibliothèque Numérique ENI',
        line_shape='spline'
    ))
    fig.add_trace(go.Scatter(
        x=pivot['Mois_str'],
        y=pivot.get('Bibliothèque Numérique ENI e-formations', []),
        mode='lines+markers',
        name='Bibliothèque Numérique ENI e-formations',
        line_shape='spline'
    ))
    fig.update_layout(
        title="Temps effectif cumulé par e-formation",
        xaxis=dict(title="Date", tickangle=-45),
        yaxis=dict(title="Heures"),
        template="plotly_white",
        margin=dict(l=60, r=60, t=100, b=100),
        legend=dict(orientation="h", x=0.5, xanchor="center", y=1.1)
    )

    st.plotly_chart(fig, use_container_width=True)

    # génération de la slide
    img_buf = BytesIO()
    fig.write_image(img_buf, format="png", width=1200, height=600, scale=2)
    img_buf.seek(0)

    prs = Presentation("TemplateGraph.pptx")
    slide = prs.slides[0]
    # calcul du centrage
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    w = Inches(9)
    h = Inches(5)
    left = (slide_w - w) // 2
    top = (slide_h - h) // 2
    slide.shapes.add_picture(img_buf, left, top, width=w, height=h)

    out_buf = BytesIO()
    prs.save(out_buf)
    out_buf.seek(0)

    st.success("Votre diapositive est prête ! 👇")
    st.download_button(
        "Télécharger la slide PPTX",
        data=out_buf.getvalue(),
        file_name="Présentation_Générée.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
